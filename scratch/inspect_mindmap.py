import collections 
import collections.abc
from pptx import Presentation
import sys

def inspect_pptx(path):
    print(f"Inspecting: {path}")
    try:
        prs = Presentation(path)
        print(f"Total slides: {len(prs.slides)}")
        for i, slide in enumerate(prs.slides):
            print(f"\n--- Slide {i+1} ---")
            text_runs = []
            image_count = 0
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text.replace("\n", " "))
                if hasattr(shape, "image"):
                    image_count += 1
            print(f"Images: {image_count}")
            print(f"Text: {' | '.join([t for t in text_runs if t.strip()])}")
    except Exception as e:
        print(f"Error: {e}")

if __name__ == "__main__":
    inspect_pptx(r"D:\tài liệu tiếng anh\lớp 1\Mind map lớp 1 Global Success.pptx")
